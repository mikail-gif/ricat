"""
PPTX Parser - Extracts slides, images, text, and tables from PPTX files.

Strategy:
  - Text & tables: parsed with python-pptx (structured API)
  - Images: parsed at raw ZIP level by scanning every r:embed / r:id
    attribute in the slide XML and mapping them through the .rels file
    to ppt/media/* entries.  This catches ALL image types regardless of
    shape kind: pictures, OLE previews, chart caches, image fills, etc.
"""

import os
import json
import hashlib
import re
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from io import BytesIO

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Relationship types that point to images/charts we want to render
_IMAGE_REL_TYPES = {
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image',
    'http://schemas.microsoft.com/office/2007/relationships/image',
}
_CHART_REL_TYPES = {
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart',
    'http://schemas.microsoft.com/office/2014/relationships/chartEx',
}
_WEB_IMG_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
                 '.tif', '.svg', '.webp', '.emf', '.wmf'}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def extract_pptx_slides(pptx_path, output_dir, department):
    """
    Extract text, tables, and images from every slide in *pptx_path*.
    Images are written to *output_dir*/media/ and referenced as web paths
    under /extracted/<department>/media/.

    Returns a dict suitable for serialising to meta.json.
    """
    media_dir = os.path.join(output_dir, 'media')
    os.makedirs(media_dir, exist_ok=True)

    # ---- pass 1: structured extraction via python-pptx (text / tables) ----
    prs = Presentation(pptx_path)

    # ---- pass 2: image extraction via raw ZIP parsing ----
    # Build a map: slide_zip_entry  ->  list of web-accessible image paths
    slide_image_map = _extract_all_slide_images(
        pptx_path, media_dir, department
    )

    slides_data = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Determine the canonical ZIP entry name for this slide
        slide_part_name = slide.part.partname  # e.g. /ppt/slides/slide1.xml
        # normalise to the key format used in slide_image_map
        slide_key = slide_part_name.lstrip('/')  # e.g. ppt/slides/slide1.xml

        slide_info = {
            "number": slide_num,
            "title": "",
            "texts": [],
            "images": slide_image_map.get(slide_key, []),
            "tables": []
        }

        # Walk shapes for text and tables
        _collect_text_and_tables(slide.shapes, slide_info)

        slides_data.append(slide_info)

    meta = {
        "totalSlides": len(slides_data),
        "extractedAt": datetime.now().isoformat(),
        "originalFile": os.path.basename(pptx_path),
        "slides": slides_data,
    }

    meta_path = os.path.join(output_dir, 'meta.json')
    with open(meta_path, 'w', encoding='utf-8') as f:
        json.dump(meta, f, indent=2, ensure_ascii=False)

    return meta


# ---------------------------------------------------------------------------
# ZIP-level image extraction
# ---------------------------------------------------------------------------

def _extract_all_slide_images(pptx_path, media_dir, department):
    """
    Open the PPTX as a ZIP and, for every slideN.xml:
      1. Parse the matching _rels/slideN.xml.rels to build rId -> target map
      2. Find every r:embed AND r:id attribute in the slide XML
      3. Resolve each rId that points to a media file
      4. For chart rIds: further look inside the chart's own .rels for images
      5. Save media files, return {slide_zip_key: [web_path, ...]}

    Deduplication is done per-slide (same bytes shown once) and globally
    (reuse the same saved file for identical images across slides).
    """
    result = {}           # slide_key -> [web_path, ...]
    seen_hashes = {}      # md5 -> web_path  (global dedup to save disk space)

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        all_names = set(zf.namelist())

        # Discover all slides from presentation.xml rels
        slide_entries = sorted(
            [n for n in all_names if re.match(r'ppt/slides/slide\d+\.xml$', n)],
            key=lambda n: int(re.search(r'(\d+)', os.path.basename(n)).group(1))
        )

        for slide_entry in slide_entries:
            slide_key = slide_entry  # e.g. "ppt/slides/slide3.xml"
            images_for_slide = []
            slide_image_hashes = set()

            # ---- Load the slide's relationship file ----
            rels_entry = _rels_path_for(slide_entry)
            rId_to_target = {}  # rId -> (zip_entry, rel_type)
            if rels_entry in all_names:
                rels_xml = zf.read(rels_entry)
                rId_to_target = _parse_rels(rels_xml, os.path.dirname(slide_entry))

            # ---- Collect every relationship ID referenced in the slide XML ----
            slide_xml = zf.read(slide_entry)
            referenced_rids = _find_all_rids(slide_xml)

            def _add(zip_entry):
                """Resolve a ZIP entry to a web path, dedup & append."""
                if zip_entry not in all_names:
                    return
                img_bytes = zf.read(zip_entry)
                h = hashlib.md5(img_bytes).hexdigest()
                if h in slide_image_hashes:
                    return
                slide_image_hashes.add(h)

                if h in seen_hashes:
                    images_for_slide.append(seen_hashes[h])
                    return

                ext = os.path.splitext(zip_entry)[1].lower()
                if ext not in _WEB_IMG_EXTS:
                    return  # not a displayable image

                slide_num = int(
                    re.search(r'(\d+)', os.path.basename(slide_entry)).group(1)
                )
                idx = len(images_for_slide) + 1
                out_ext = ext.lstrip('.')
                if out_ext in ('jpeg',):
                    out_ext = 'jpg'

                # Convert EMF/WMF → PNG when Pillow is available
                if ext in ('.emf', '.wmf') and HAS_PIL:
                    try:
                        img = Image.open(BytesIO(img_bytes))
                        out_ext = 'png'
                        filename = f"s{slide_num}_i{idx}.png"
                        out_path = os.path.join(media_dir, filename)
                        img.save(out_path, 'PNG')
                    except Exception:
                        filename = f"s{slide_num}_i{idx}.{out_ext}"
                        out_path = os.path.join(media_dir, filename)
                        with open(out_path, 'wb') as f:
                            f.write(img_bytes)
                else:
                    filename = f"s{slide_num}_i{idx}.{out_ext}"
                    out_path = os.path.join(media_dir, filename)
                    with open(out_path, 'wb') as f:
                        f.write(img_bytes)

                web_path = f"/extracted/{department}/media/{filename}"
                seen_hashes[h] = web_path
                images_for_slide.append(web_path)

            # ---- Process each referenced rId ----
            for rId in referenced_rids:
                if rId not in rId_to_target:
                    continue
                target_entry, rel_type = rId_to_target[rId]

                if rel_type in _IMAGE_REL_TYPES:
                    # Direct image reference
                    _add(target_entry)

                elif rel_type in _CHART_REL_TYPES:
                    # Chart: look inside chart's own .rels for cached images
                    chart_rels_entry = _rels_path_for(target_entry)
                    if chart_rels_entry in all_names:
                        chart_rels_xml = zf.read(chart_rels_entry)
                        chart_rId_map = _parse_rels(
                            chart_rels_xml, os.path.dirname(target_entry)
                        )
                        for _, (chart_target, chart_rel_type) in chart_rId_map.items():
                            if chart_rel_type in _IMAGE_REL_TYPES:
                                _add(chart_target)

            # ---- Catch any remaining ppt/media/* files referenced by the
            #      slide but whose rId wasn't in the collected set
            #      (some tools emit non-standard attribute names).
            # We do this by scanning the slide XML for media file names. ----
            for rId, (target_entry, rel_type) in rId_to_target.items():
                if rel_type in _IMAGE_REL_TYPES:
                    _add(target_entry)

            result[slide_key] = images_for_slide

    return result


def _rels_path_for(part_path):
    """Return the .rels path for a given part path inside the ZIP."""
    dirname = os.path.dirname(part_path)
    basename = os.path.basename(part_path)
    return f"{dirname}/_rels/{basename}.rels"


def _parse_rels(rels_xml_bytes, base_dir):
    """
    Parse a .rels XML and return {rId: (resolved_zip_entry, rel_type)}.
    """
    result = {}
    try:
        root = ET.fromstring(rels_xml_bytes)
    except ET.ParseError:
        return result

    ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
    for rel in root.findall(f'{{{ns}}}Relationship'):
        rId = rel.get('Id', '')
        rel_type = rel.get('Type', '')
        target = rel.get('Target', '')
        target_mode = rel.get('TargetMode', 'Internal')

        if target_mode == 'External' or not target:
            continue

        # Resolve relative path
        if target.startswith('/'):
            resolved = target.lstrip('/')
        else:
            resolved = os.path.normpath(
                os.path.join(base_dir, target)
            ).replace('\\', '/')

        result[rId] = (resolved, rel_type)
    return result


def _find_all_rids(xml_bytes):
    """
    Return every unique r:embed / r:id / r:link value found anywhere
    in the XML bytes.  We use regex on raw bytes for speed and to avoid
    namespace headaches with uncommon OOXML writers.
    """
    # Match r:embed="rId..." or r:id="rId..." or r:link="rId..."
    pattern = rb'r:(?:embed|id|link)="(rId\d+)"'
    return set(m.group(1).decode() for m in re.finditer(pattern, xml_bytes))


# ---------------------------------------------------------------------------
# python-pptx pass: text and tables only
# ---------------------------------------------------------------------------

def _collect_text_and_tables(shapes, slide_info):
    """Walk all shapes recursively to collect text paragraphs and tables."""
    for shape in shapes:
        # Text
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_info["texts"].append(text)
                    if not slide_info["title"]:
                        slide_info["title"] = text

        # Group – recurse
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_text_and_tables(shape.shapes, slide_info)

        # Table
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])
            if table_data:
                slide_info["tables"].append(table_data)
